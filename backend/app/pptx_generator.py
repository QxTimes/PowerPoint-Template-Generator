from pathlib import Path
from tempfile import NamedTemporaryFile
from pptx import Presentation


def create_pptx(prompt: str) -> str:
    """Generate a PowerPoint presentation with the prompt text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = prompt

    tmp = NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name
